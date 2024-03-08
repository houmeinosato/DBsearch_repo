import streamlit as st
from fpdf import FPDF
import base64

pdf = FPDF()  # pdf object
pdf = FPDF(orientation="P", unit="mm", format="A4")
pdf.add_page()

pdf.set_font("Times", "B", 18)
pdf.set_xy(10.0, 20)
pdf.cell(w=75.0, h=5.0, align="L", txt="This is my sample text")

st.download_button(
    "Download Report",
    data=pdf.output(dest='S').encode('latin-1'),
    file_name="Output.pdf",
)


'''
import streamlit as st
import base64
from streamlit.components.v1 import html
import numpy as np
from PIL import Image
import io

def create_image_link():
    img_array = np.zeros((256, 256, 3), np.uint8)
    img_array[:, :, 0] = 255
    img_array[:, :, 1] = np.arange(256, dtype=np.uint8)[None, :]
    img_array[:, :, 2] = np.arange(256, dtype=np.uint8)[:, None]
    with io.BytesIO() as buf:
        with Image.fromarray(img_array) as img:
            img.save(buf, format="png")
        image_str = base64.b64encode(buf.getvalue()).decode()
        js_code = f"""<a href="data:png;base64,{image_str}" download="sample.png">download</a>"""
    return js_code

def main():
    download_button = st.button("Click to download")
    container = st.empty()
    if download_button:
        with container:
            html(create_image_link(), height=50)

if __name__ == "__main__":
    main()
'''



'''
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO


# パワーポイントファイルを作成する関数
def create_ppt(text, images):
    prs = Presentation()
    # スライドのサイズを16:9に設定
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 新しいスライドを追加（空白のスライドレイアウトを使用）
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # テキストをスライドに追加（テキストがある場合）
    if text:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(2))
        p = textbox.text_frame.add_paragraph()
        p.text = text
        # フォントサイズを設定（テキスト量に応じて調整が必要かもしれません）
        p.font.size = Pt(18)

    # 画像をスライドに追加（画像がある場合）
    for idx, image in enumerate(images):
        if image:
            # BytesIOオブジェクトを作成し、Streamlitからアップロードされた画像データを読み込む
            image_stream = BytesIO(image.getvalue())
            left = Inches(0.5) + (Inches(5) * idx)  # 画像を横に並べる
            top = Inches(3)  # テキストボックスの下
            height = Inches(5.5)  # 画像の高さを適切に設定
            slide.shapes.add_picture(image_stream, left, top, height=height)

    # バイナリストリームに保存
    bio = BytesIO()
    prs.save(bio)
    return bio


# Streamlit UI
st.title('パワーポイントファイル生成アプリ')

# テキスト入力
user_input_text = st.text_area("スライドに挿入するテキストを入力してください")

# 画像アップローダー
uploaded_images = []
for i in range(3):  # 最大3枚の画像
    uploaded_image = st.file_uploader(f"画像{i+1}をアップロードしてください（任意）", 
                                      key=f"image_{i}", 
                                      type=['jpg', 'jpeg', 'png']
                                      )
    if uploaded_image:
        uploaded_images.append(uploaded_image)

# 生成ボタン
if st.button('パワーポイントファイルを生成'):
    if not user_input_text and not uploaded_images:
        st.warning('テキストまたは画像をアップロードしてください。')
    else:
        # パワーポイントファイルを作成
        pptx_bio = create_ppt(user_input_text, uploaded_images)
        pptx_bio.seek(0)

        # ファイルをダウンロードするためのリンクを生成
        st.download_button(label='ダウンロード', 
                           data=pptx_bio, 
                           file_name='custom_presentation.pptx', 
                           mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                           )
'''
